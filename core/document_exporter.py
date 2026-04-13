"""
ドキュメントエクスポーター (Document Exporter)

アイが調べた情報や会話内容を Office ドキュメントに出力する。

対応フォーマット:
  - Word (.docx)   — レポート・資料・議事録
  - PowerPoint (.pptx) — プレゼン資料・まとめスライド
  - Excel (.xlsx)  — データ一覧・分析結果・比較表

使い方:
  「Wordにまとめて」「パワポで資料作って」「エクセルに一覧作って」
"""
from __future__ import annotations

import logging
import re
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# オプショナルインポート（未インストールでも起動可能）
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

try:
    from docx import Document as DocxDocument
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# データ構造
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

@dataclass(frozen=True)
class DocumentSection:
    """ドキュメントの1セクション"""
    heading: str
    body: str = ""
    bullet_points: tuple[str, ...] = ()
    table_headers: tuple[str, ...] = ()
    table_rows: tuple[tuple[str, ...], ...] = ()


@dataclass(frozen=True)
class DocumentContent:
    """エクスポート用のドキュメント内容"""
    title: str
    subtitle: str = ""
    sections: tuple[DocumentSection, ...] = ()
    author: str = "アイ"
    created_at: float = 0.0


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# パーサー: テキスト → DocumentContent
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class ContentParser:
    """
    マークダウン風テキストを DocumentContent に変換する。

    入力フォーマット:
      # タイトル
      ## セクション1
      本文テキスト
      - 箇条書き1
      - 箇条書き2
      | ヘッダ1 | ヘッダ2 |
      | データ1 | データ2 |
    """

    _HEADING_RE = re.compile(r"^(#{1,3})\s+(.+)")
    _BULLET_RE = re.compile(r"^[-*]\s+(.+)")
    _TABLE_RE = re.compile(r"^\|(.+)\|$")
    _TABLE_SEP_RE = re.compile(r"^\|[-\s|:]+\|$")

    def parse(self, text: str) -> DocumentContent:
        """テキストを解析して DocumentContent を生成"""
        lines = text.strip().splitlines()
        title = ""
        subtitle = ""
        sections: list[DocumentSection] = []
        current_heading = ""
        current_body_lines: list[str] = []
        current_bullets: list[str] = []
        current_table_headers: list[str] = []
        current_table_rows: list[tuple[str, ...]] = []

        def _flush_section() -> None:
            nonlocal current_heading, current_body_lines, current_bullets
            nonlocal current_table_headers, current_table_rows
            if current_heading:
                sections.append(DocumentSection(
                    heading=current_heading,
                    body="\n".join(current_body_lines).strip(),
                    bullet_points=tuple(current_bullets),
                    table_headers=tuple(current_table_headers),
                    table_rows=tuple(
                        tuple(r) for r in current_table_rows
                    ),
                ))
            current_heading = ""
            current_body_lines = []
            current_bullets = []
            current_table_headers = []
            current_table_rows = []

        for line in lines:
            line_stripped = line.strip()

            # 見出し
            hm = self._HEADING_RE.match(line_stripped)
            if hm:
                level = len(hm.group(1))
                text_content = hm.group(2).strip()
                if level == 1 and not title:
                    title = text_content
                elif level == 1 and not subtitle:
                    subtitle = text_content
                else:
                    _flush_section()
                    current_heading = text_content
                continue

            # テーブル区切り行（スキップ）
            if self._TABLE_SEP_RE.match(line_stripped):
                continue

            # テーブル行
            tm = self._TABLE_RE.match(line_stripped)
            if tm:
                cells = [c.strip() for c in tm.group(1).split("|")]
                if not current_table_headers:
                    current_table_headers = cells
                else:
                    current_table_rows.append(tuple(cells))
                continue

            # 箇条書き
            bm = self._BULLET_RE.match(line_stripped)
            if bm:
                current_bullets.append(bm.group(1))
                continue

            # 本文
            if line_stripped:
                current_body_lines.append(line_stripped)

        _flush_section()

        if not title:
            title = "ドキュメント"

        return DocumentContent(
            title=title,
            subtitle=subtitle,
            sections=tuple(sections),
            created_at=time.time(),
        )


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Word エクスポーター
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class WordExporter:
    """DocumentContent → .docx"""

    def export(self, content: DocumentContent, output_path: Path) -> Path:
        if not HAS_DOCX:
            raise RuntimeError(
                "python-docx が未インストールです。"
                "pip install python-docx を実行してください。"
            )
        doc = DocxDocument()

        # タイトル
        doc.add_heading(content.title, level=0)
        if content.subtitle:
            sub = doc.add_paragraph(content.subtitle)
            sub.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # セクション
        for section in content.sections:
            doc.add_heading(section.heading, level=1)

            if section.body:
                doc.add_paragraph(section.body)

            for bullet in section.bullet_points:
                doc.add_paragraph(bullet, style="List Bullet")

            if section.table_headers:
                table = doc.add_table(
                    rows=1, cols=len(section.table_headers)
                )
                table.style = "Table Grid"
                for i, header in enumerate(section.table_headers):
                    cell = table.rows[0].cells[i]
                    cell.text = header
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

                for row_data in section.table_rows:
                    row = table.add_row().cells
                    for i, val in enumerate(row_data):
                        if i < len(row):
                            row[i].text = val

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
        logger.info("Word ドキュメント生成: %s", output_path)
        return output_path


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# PowerPoint エクスポーター
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class PowerPointExporter:
    """DocumentContent → .pptx"""

    def export(self, content: DocumentContent, output_path: Path) -> Path:
        if not HAS_PPTX:
            raise RuntimeError(
                "python-pptx が未インストールです。"
                "pip install python-pptx を実行してください。"
            )
        prs = Presentation()

        # タイトルスライド
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = content.title
        if content.subtitle and slide.placeholders[1]:
            slide.placeholders[1].text = content.subtitle

        # セクションごとにスライド生成
        for section in content.sections:
            # テーブルがあるセクション → 空白スライド + テーブル
            if section.table_headers:
                self._add_table_slide(prs, section)
            else:
                self._add_content_slide(prs, section)

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        logger.info("PowerPoint 生成: %s", output_path)
        return output_path

    def _add_content_slide(
        self, prs: Any, section: DocumentSection
    ) -> None:
        """コンテンツスライドを追加"""
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = section.heading

        body = ""
        if section.body:
            body += section.body + "\n"
        for bullet in section.bullet_points:
            body += f"• {bullet}\n"

        if body and slide.placeholders[1]:
            slide.placeholders[1].text = body.strip()

    def _add_table_slide(
        self, prs: Any, section: DocumentSection
    ) -> None:
        """テーブル付きスライドを追加"""
        # タイトル付き空白
        layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = section.heading

        rows = len(section.table_rows) + 1
        cols = len(section.table_headers)
        left = PptxInches(0.5)
        top = PptxInches(1.8)
        width = PptxInches(9.0)
        height = PptxInches(0.4 * rows)

        table = slide.shapes.add_table(
            rows, cols, left, top, width, height
        ).table

        for col_idx, header in enumerate(section.table_headers):
            table.cell(0, col_idx).text = header

        for row_idx, row_data in enumerate(section.table_rows, start=1):
            for col_idx, val in enumerate(row_data):
                if col_idx < cols:
                    table.cell(row_idx, col_idx).text = val


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# Excel エクスポーター
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class ExcelExporter:
    """DocumentContent → .xlsx"""

    _HEADER_FILL = PatternFill(
        fill_type="solid", fgColor="4472C4"
    ) if HAS_OPENPYXL else None
    _HEADER_FONT = Font(
        bold=True, color="FFFFFF", size=11
    ) if HAS_OPENPYXL else None

    def export(self, content: DocumentContent, output_path: Path) -> Path:
        if not HAS_OPENPYXL:
            raise RuntimeError(
                "openpyxl が未インストールです。"
                "pip install openpyxl を実行してください。"
            )
        wb = Workbook()
        ws = wb.active
        ws.title = content.title[:31]  # Excel sheet name max 31 chars

        row_num = 1

        # タイトル
        ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=5)
        title_cell = ws.cell(row=1, column=1, value=content.title)
        title_cell.font = Font(bold=True, size=16)
        title_cell.alignment = Alignment(horizontal="center")
        row_num += 1

        if content.subtitle:
            ws.merge_cells(
                start_row=2, start_column=1, end_row=2, end_column=5
            )
            sub_cell = ws.cell(row=2, column=1, value=content.subtitle)
            sub_cell.font = Font(size=12, italic=True)
            sub_cell.alignment = Alignment(horizontal="center")
            row_num += 1

        row_num += 1  # 空行

        for section in content.sections:
            # セクション見出し
            ws.cell(
                row=row_num, column=1, value=section.heading
            ).font = Font(bold=True, size=13)
            row_num += 1

            # 本文
            if section.body:
                for body_line in section.body.splitlines():
                    ws.cell(row=row_num, column=1, value=body_line)
                    row_num += 1

            # 箇条書き
            for bullet in section.bullet_points:
                ws.cell(row=row_num, column=1, value=f"• {bullet}")
                row_num += 1

            # テーブル
            if section.table_headers:
                for col_idx, header in enumerate(
                    section.table_headers, start=1
                ):
                    cell = ws.cell(
                        row=row_num, column=col_idx, value=header
                    )
                    cell.font = self._HEADER_FONT
                    cell.fill = self._HEADER_FILL
                    cell.alignment = Alignment(horizontal="center")
                row_num += 1

                for row_data in section.table_rows:
                    for col_idx, val in enumerate(row_data, start=1):
                        ws.cell(row=row_num, column=col_idx, value=val)
                    row_num += 1

            row_num += 1  # セクション間の空行

        # 列幅自動調整（概算）
        for col_idx in range(1, 8):
            ws.column_dimensions[
                get_column_letter(col_idx)
            ].width = 18

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(output_path))
        logger.info("Excel ファイル生成: %s", output_path)
        return output_path


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# 統合エンジン
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

class DocumentExportEngine:
    """
    ドキュメントエクスポートの統合エンジン。

    使い方:
      engine = DocumentExportEngine(output_dir)
      path = engine.export_word("# タイトル\n## セクション\n本文")
      path = engine.export_pptx("# タイトル\n...")
      path = engine.export_excel("# タイトル\n...")
    """

    def __init__(self, output_dir: str | Path) -> None:
        self._output_dir = Path(output_dir)
        self._output_dir.mkdir(parents=True, exist_ok=True)
        self._parser = ContentParser()
        self._word = WordExporter()
        self._pptx = PowerPointExporter()
        self._excel = ExcelExporter()

    def _make_filename(self, title: str, ext: str) -> Path:
        """タイトルからファイル名を生成"""
        safe = re.sub(r'[\\/:*?"<>|]', "_", title)[:50]
        ts = int(time.time())
        return self._output_dir / f"{safe}_{ts}.{ext}"

    def export_word(
        self,
        text: str,
        filename: str | None = None,
    ) -> Path:
        """テキスト → Word ドキュメント"""
        content = self._parser.parse(text)
        path = (
            self._output_dir / filename
            if filename
            else self._make_filename(content.title, "docx")
        )
        return self._word.export(content, path)

    def export_pptx(
        self,
        text: str,
        filename: str | None = None,
    ) -> Path:
        """テキスト → PowerPoint"""
        content = self._parser.parse(text)
        path = (
            self._output_dir / filename
            if filename
            else self._make_filename(content.title, "pptx")
        )
        return self._pptx.export(content, path)

    def export_excel(
        self,
        text: str,
        filename: str | None = None,
    ) -> Path:
        """テキスト → Excel"""
        content = self._parser.parse(text)
        path = (
            self._output_dir / filename
            if filename
            else self._make_filename(content.title, "xlsx")
        )
        return self._excel.export(content, path)

    def get_available_formats(self) -> list[str]:
        """利用可能なフォーマットを返す"""
        formats = []
        if HAS_DOCX:
            formats.append("word")
        if HAS_PPTX:
            formats.append("pptx")
        if HAS_OPENPYXL:
            formats.append("excel")
        return formats

    def get_status_text(self) -> str:
        """ステータステキスト"""
        fmts = self.get_available_formats()
        if not fmts:
            return "📄 ドキュメント出力: 未インストール"
        labels = {
            "word": "Word(.docx)",
            "pptx": "PowerPoint(.pptx)",
            "excel": "Excel(.xlsx)",
        }
        available = "、".join(labels.get(f, f) for f in fmts)
        return f"📄 ドキュメント出力: {available} 対応"
