"""
Tests for core.document_exporter module.

Covers ContentParser, WordExporter, PowerPointExporter, ExcelExporter,
DocumentExportEngine, and frozen dataclasses (DocumentContent, DocumentSection).
"""
from __future__ import annotations

import tempfile
from dataclasses import FrozenInstanceError
from pathlib import Path

import pytest

from core.document_exporter import (
    ContentParser,
    DocumentContent,
    DocumentExportEngine,
    DocumentSection,
    ExcelExporter,
    PowerPointExporter,
    WordExporter,
)


# ──────────────────────────────────────────────
# Fixtures
# ──────────────────────────────────────────────

SAMPLE_MARKDOWN = """\
# テスト資料
# サブタイトル
## セクション1
本文テキストです。
- 箇条書き1
- 箇条書き2
## セクション2
| 名前 | 年齢 |
|------|------|
| 太郎 | 25 |
| 花子 | 30 |
"""

MINIMAL_MARKDOWN = """\
# Hello
## Section A
Body text here.
"""


@pytest.fixture()
def parser() -> ContentParser:
    return ContentParser()


@pytest.fixture()
def tmp_dir() -> str:
    return tempfile.mkdtemp()


@pytest.fixture()
def sample_content(parser: ContentParser) -> DocumentContent:
    return parser.parse(SAMPLE_MARKDOWN)


@pytest.fixture()
def minimal_content(parser: ContentParser) -> DocumentContent:
    return parser.parse(MINIMAL_MARKDOWN)


# ──────────────────────────────────────────────
# DocumentSection / DocumentContent frozen tests
# ──────────────────────────────────────────────

class TestDocumentSectionFrozen:
    def test_frozen_heading(self) -> None:
        section = DocumentSection(heading="Test")
        with pytest.raises(FrozenInstanceError):
            section.heading = "Changed"  # type: ignore[misc]

    def test_defaults(self) -> None:
        section = DocumentSection(heading="H")
        assert section.body == ""
        assert section.bullet_points == ()
        assert section.table_headers == ()
        assert section.table_rows == ()


class TestDocumentContentFrozen:
    def test_frozen_title(self) -> None:
        content = DocumentContent(title="T")
        with pytest.raises(FrozenInstanceError):
            content.title = "Changed"  # type: ignore[misc]

    def test_defaults(self) -> None:
        content = DocumentContent(title="T")
        assert content.subtitle == ""
        assert content.sections == ()
        assert content.author == "アイ"


# ──────────────────────────────────────────────
# ContentParser tests
# ──────────────────────────────────────────────

class TestContentParser:
    def test_title_extraction(self, parser: ContentParser) -> None:
        content = parser.parse("# My Title\n## Sec\nBody")
        assert content.title == "My Title"

    def test_subtitle_extraction(self, parser: ContentParser) -> None:
        content = parser.parse(SAMPLE_MARKDOWN)
        assert content.subtitle == "サブタイトル"

    def test_section_count(self, sample_content: DocumentContent) -> None:
        assert len(sample_content.sections) == 2

    def test_section_heading(self, sample_content: DocumentContent) -> None:
        assert sample_content.sections[0].heading == "セクション1"

    def test_body_text(self, sample_content: DocumentContent) -> None:
        assert "本文テキスト" in sample_content.sections[0].body

    def test_bullet_points(self, sample_content: DocumentContent) -> None:
        bullets = sample_content.sections[0].bullet_points
        assert len(bullets) == 2
        assert bullets[0] == "箇条書き1"
        assert bullets[1] == "箇条書き2"

    def test_table_headers(self, sample_content: DocumentContent) -> None:
        headers = sample_content.sections[1].table_headers
        assert "名前" in headers
        assert "年齢" in headers

    def test_table_rows(self, sample_content: DocumentContent) -> None:
        rows = sample_content.sections[1].table_rows
        assert len(rows) == 2
        assert "太郎" in rows[0]
        assert "30" in rows[1]

    def test_empty_input(self, parser: ContentParser) -> None:
        content = parser.parse("")
        assert content.title == "ドキュメント"
        assert content.sections == ()

    def test_no_title_defaults(self, parser: ContentParser) -> None:
        content = parser.parse("just some text without headings")
        assert content.title == "ドキュメント"

    def test_created_at_set(self, parser: ContentParser) -> None:
        content = parser.parse("# T\n## S\nBody")
        assert content.created_at > 0


# ──────────────────────────────────────────────
# WordExporter tests
# ──────────────────────────────────────────────

class TestWordExporter:
    def test_creates_file(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        path = Path(tmp_dir) / "test.docx"
        result = WordExporter().export(sample_content, path)
        assert result.exists()
        assert result.stat().st_size > 0

    def test_file_is_valid_docx(
        self, tmp_dir: str, minimal_content: DocumentContent
    ) -> None:
        from docx import Document as DocxDocument

        path = Path(tmp_dir) / "valid.docx"
        WordExporter().export(minimal_content, path)
        doc = DocxDocument(str(path))
        # Title should be present (level=0 uses "Title" style, not "Heading")
        all_text = [p.text for p in doc.paragraphs]
        assert any("Hello" in t for t in all_text)

    def test_tables_written(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from docx import Document as DocxDocument

        path = Path(tmp_dir) / "tables.docx"
        WordExporter().export(sample_content, path)
        doc = DocxDocument(str(path))
        assert len(doc.tables) >= 1

    def test_bullet_points_written(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from docx import Document as DocxDocument

        path = Path(tmp_dir) / "bullets.docx"
        WordExporter().export(sample_content, path)
        doc = DocxDocument(str(path))
        bullet_texts = [
            p.text for p in doc.paragraphs if p.style.name == "List Bullet"
        ]
        assert len(bullet_texts) == 2


# ──────────────────────────────────────────────
# PowerPointExporter tests
# ──────────────────────────────────────────────

class TestPowerPointExporter:
    def test_creates_file(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        path = Path(tmp_dir) / "test.pptx"
        result = PowerPointExporter().export(sample_content, path)
        assert result.exists()
        assert result.stat().st_size > 0

    def test_file_is_valid_pptx(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from pptx import Presentation

        path = Path(tmp_dir) / "valid.pptx"
        PowerPointExporter().export(sample_content, path)
        prs = Presentation(str(path))
        # Title slide + 2 section slides
        assert len(prs.slides) >= 3

    def test_title_slide_text(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from pptx import Presentation

        path = Path(tmp_dir) / "title.pptx"
        PowerPointExporter().export(sample_content, path)
        prs = Presentation(str(path))
        title_slide = prs.slides[0]
        assert title_slide.shapes.title.text == "テスト資料"

    def test_table_slide_has_table(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from pptx import Presentation

        path = Path(tmp_dir) / "table_slide.pptx"
        PowerPointExporter().export(sample_content, path)
        prs = Presentation(str(path))
        # The last slide should be the table slide (section 2)
        table_slide = prs.slides[-1]
        has_table = any(shape.has_table for shape in table_slide.shapes)
        assert has_table


# ──────────────────────────────────────────────
# ExcelExporter tests
# ──────────────────────────────────────────────

class TestExcelExporter:
    def test_creates_file(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        path = Path(tmp_dir) / "test.xlsx"
        result = ExcelExporter().export(sample_content, path)
        assert result.exists()
        assert result.stat().st_size > 0

    def test_file_is_valid_xlsx(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from openpyxl import load_workbook

        path = Path(tmp_dir) / "valid.xlsx"
        ExcelExporter().export(sample_content, path)
        wb = load_workbook(str(path))
        ws = wb.active
        assert ws is not None

    def test_title_in_first_cell(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from openpyxl import load_workbook

        path = Path(tmp_dir) / "title.xlsx"
        ExcelExporter().export(sample_content, path)
        wb = load_workbook(str(path))
        ws = wb.active
        assert ws.cell(row=1, column=1).value == "テスト資料"

    def test_header_styling(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from openpyxl import load_workbook

        path = Path(tmp_dir) / "style.xlsx"
        ExcelExporter().export(sample_content, path)
        wb = load_workbook(str(path))
        ws = wb.active
        title_cell = ws.cell(row=1, column=1)
        assert title_cell.font.bold is True

    def test_data_rows_present(
        self, tmp_dir: str, sample_content: DocumentContent
    ) -> None:
        from openpyxl import load_workbook

        path = Path(tmp_dir) / "data.xlsx"
        ExcelExporter().export(sample_content, path)
        wb = load_workbook(str(path))
        ws = wb.active
        # Scan all cells for a known data value
        found = False
        for row in ws.iter_rows(values_only=True):
            if "太郎" in (str(c) for c in row if c is not None):
                found = True
                break
        assert found


# ──────────────────────────────────────────────
# DocumentExportEngine tests
# ──────────────────────────────────────────────

class TestDocumentExportEngine:
    def test_export_word(self, tmp_dir: str) -> None:
        engine = DocumentExportEngine(tmp_dir)
        path = engine.export_word(SAMPLE_MARKDOWN)
        assert path.exists()
        assert path.suffix == ".docx"
        assert path.stat().st_size > 0

    def test_export_pptx(self, tmp_dir: str) -> None:
        engine = DocumentExportEngine(tmp_dir)
        path = engine.export_pptx(SAMPLE_MARKDOWN)
        assert path.exists()
        assert path.suffix == ".pptx"
        assert path.stat().st_size > 0

    def test_export_excel(self, tmp_dir: str) -> None:
        engine = DocumentExportEngine(tmp_dir)
        path = engine.export_excel(SAMPLE_MARKDOWN)
        assert path.exists()
        assert path.suffix == ".xlsx"
        assert path.stat().st_size > 0

    def test_export_word_custom_filename(self, tmp_dir: str) -> None:
        engine = DocumentExportEngine(tmp_dir)
        path = engine.export_word(SAMPLE_MARKDOWN, filename="custom.docx")
        assert path.name == "custom.docx"
        assert path.exists()

    def test_filename_generation(self, tmp_dir: str) -> None:
        engine = DocumentExportEngine(tmp_dir)
        generated = engine._make_filename("テスト資料", "docx")
        assert generated.suffix == ".docx"
        assert "テスト資料" in generated.stem

    def test_available_formats(self, tmp_dir: str) -> None:
        engine = DocumentExportEngine(tmp_dir)
        formats = engine.get_available_formats()
        assert "word" in formats
        assert "pptx" in formats
        assert "excel" in formats

    def test_status_text(self, tmp_dir: str) -> None:
        engine = DocumentExportEngine(tmp_dir)
        status = engine.get_status_text()
        assert "ドキュメント出力" in status
